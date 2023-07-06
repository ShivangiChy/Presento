import collections.abc
from pptx import Presentation
import json
with open('output11.json') as f:
    data = json.load(f)
# print(data["slide"])

def _add_leveled_bullet(_placeholder, _text, level=0):
    _prg = _placeholder.text_frame.add_paragraph()
    _prg.level = level
    _prg.text = _text


prs = Presentation()
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
# subtitle = slide.placeholders[1]
title.text = data['presentationTitle']


for i in range(len(data['slide'])):
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = data["slide"][i]["title"]
    for j in range(len(data["slide"][i]["points"])):
        _add_leveled_bullet(subtitle,data["slide"][i]["points"][j] , 0)
    # _add_leveled_bullet(subtitle,data["slide"][i]["points"][9] , 0)
    # _add_leveled_bullet(subtitle,data["slide"][i]["imageDescription"] , 0)
#
# subtitle.text = "python-pptx was here!"
# _add_leveled_bullet(subtitle,"python-pptx was here!" , 0)
# subtitle.text = "python-pptx was here!"
prs.save('test5.pptx')